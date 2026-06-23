import streamlit as st
import openpyxl
import re
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.units import cm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.enums import TA_LEFT, TA_CENTER
import os, urllib.request

# ============================================================
# ШРИФТЫ (ищем в системе или используем встроенные пути)
# ============================================================
FONT_PATHS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/DejaVuSans.ttf",
]
FONTB_PATHS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/DejaVuSans-Bold.ttf",
]

def find_font(paths):
    for p in paths:
        if os.path.exists(p):
            return p
    return None

font_path  = find_font(FONT_PATHS)
fontb_path = find_font(FONTB_PATHS)

if font_path and fontb_path:
    pdfmetrics.registerFont(TTFont("DejaVu",     font_path))
    pdfmetrics.registerFont(TTFont("DejaVuBold", fontb_path))
else:
    # Устанавливаем шрифты через apt если не найдены
    import subprocess
    subprocess.run(["apt-get", "install", "-y", "-q", "fonts-dejavu-core"], capture_output=True)
    font_path  = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    fontb_path = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
    pdfmetrics.registerFont(TTFont("DejaVu",     font_path))
    pdfmetrics.registerFont(TTFont("DejaVuBold", fontb_path))

# ============================================================
# ЦВЕТА
# ============================================================
ORANGE = RGBColor(0xE8, 0x77, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
CREAM  = RGBColor(0xFF, 0xF8, 0xF0)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
RED    = RGBColor(0xDC, 0x26, 0x26)
GRAY   = RGBColor(0x88, 0x88, 0x88)

ORANGE_PDF = colors.HexColor("#E87722")
CREAM_PDF  = colors.HexColor("#FFF8F0")
GREEN_PDF  = colors.HexColor("#16A34A")
RED_PDF    = colors.HexColor("#DC2626")
GRAY_PDF   = colors.HexColor("#888888")
DARK_PDF   = colors.HexColor("#1A1A1A")

# ============================================================
# СПИСКИ СОТРУДНИКОВ
# ============================================================
DEPT_MGR  = ["Абрамович","Кузменко","Равба","Кузнецов","Вайтусенок В.М"]
CONS_LIST = ["Гарькавенко","Ляшук","Пальчех","Анищик"]
OTHER_MGR = ["Афанасович","Горбач","Каледа","Мандрик","Марко","Чепляев","Вайтусенок М.В"]

def match_list(name, lst):
    for m in lst:
        parts = m.split()
        if parts[0] not in name: continue
        if len(parts) > 1:
            if parts[1] in name: return True
        else:
            return True
    return False

# ============================================================
# ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ
# ============================================================
def fmt(n, dec=0):
    try:
        v = float(n)
        s = f"{v:,.{dec}f}".replace(",", " ")
        return s
    except:
        return "—"

def pct(cur, prev):
    try:
        c, p = float(cur or 0), float(prev or 0)
        if p == 0: return None
        return round(((c - p) / abs(p)) * 100, 1)
    except:
        return None

def sign(v):
    return "+" if v and v > 0 else ""

def dstr(v):
    if v is None: return "—"
    return f"{sign(v)}{v}%"

def dcol_pptx(v):
    if v is None: return GRAY
    return GREEN if v > 0 else RED if v < 0 else GRAY

def dcol_pdf(v):
    if v is None: return GRAY_PDF
    return GREEN_PDF if v > 0 else RED_PDF if v < 0 else GRAY_PDF

# ============================================================
# ПАРСИНГ EXCEL
# ============================================================
def parse_excel(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    all_rows = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            all_rows.append(list(row))

    def find_by_label(kw):
        for row in all_rows:
            lbl = str(row[0] if row[0] else (row[1] if len(row)>1 else "")).lower()
            if kw.lower() in lbl:
                nums = sorted([float(c) for c in row if isinstance(c,(int,float)) and c > 0])
                if len(nums) >= 3:
                    return {"vol": nums[1], "revenue": nums[-1], "gp": nums[-2], "margin": nums[0]}
        return {}

    premixes     = find_by_label("премикс")
    compound     = find_by_label("комбикорм")
    additives    = find_by_label("всего кормовые добавки")
    concentrates = find_by_label("кормовых концентрат")

    dept_managers, other_managers, consultants = [], [], []
    name_re = re.compile(r'^[А-ЯЁ][а-яёА-ЯЁ\-]+\s+[А-ЯЁ]\.[А-ЯЁ]')
    skip = ["всего","итого","менеджер","показатели","информация"]

    for row in all_rows:
        name = str(row[0] or "").strip()
        if not name_re.match(name): continue
        if any(name.lower().startswith(s) for s in skip): continue
        nums = [float(c) for c in row[1:] if isinstance(c,(int,float)) and c > 0]
        if not nums: continue
        gp = nums[0] if len(nums)==1 else (nums[0]-nums[1] if len(nums)==2 else nums[-1])
        if gp < 100: continue
        if match_list(name, DEPT_MGR):
            if not any(x["name"]==name for x in dept_managers):
                dept_managers.append({"name":name,"gp":gp,"gp_prev":0})
        elif match_list(name, CONS_LIST):
            if not any(x["name"]==name for x in consultants):
                consultants.append({"name":name,"gp":gp,"gp_prev":0})
        elif match_list(name, OTHER_MGR):
            if not any(x["name"]==name for x in other_managers):
                other_managers.append({"name":name,"gp":gp,"gp_prev":0})

    dept = {"gp_arrival":0,"transport":0,"gp_salary":0}
    dir_total = {"gp":0,"transport":0}
    for row in all_rows:
        lbl = str(row[0] or "").lower()
        nums = [float(c) for c in row if isinstance(c,(int,float)) and c > 0]
        if "всего по" in lbl and "кормлени" in lbl:
            if len(nums)>=2: dir_total["gp"]=nums[0]; dir_total["transport"]=nums[1]
        if (lbl=="всего" or lbl.startswith("всего ")) and "кормлени" not in lbl and "добавки" not in lbl and "концентрат" not in lbl:
            big = [n for n in nums if n>1000]
            if len(big)>=2 and dept["gp_arrival"]==0:
                dept["gp_arrival"]=big[-1]
                tr = next((n for n in big if 1000<n<30000),0)
                dept["transport"]=tr
                dept["gp_salary"]=dept["gp_arrival"]-tr

    return {"premixes":premixes,"compound":compound,"additives":additives,
            "concentrates":concentrates,"dept_managers":dept_managers,
            "other_managers":other_managers,"consultants":consultants,
            "dept":dept,"dir_total":dir_total}

# ============================================================
# ГЕНЕРАЦИЯ PPTX
# ============================================================
def add_rect(slide, x,y,w,h, color):
    s = slide.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=color; s.line.fill.background()
    return s

def add_text(slide,text,x,y,w,h,size=12,bold=False,color=DARK,align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.alignment=align
    r = p.add_run(); r.text=str(text)
    r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return tb

def build_pptx(nd, pd, month, prev_month):
    prs = Presentation()
    prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
    W,H = 13.33,7.5

    def blank(bg=CREAM):
        sl=prs.slide_layouts[6]; slide=prs.slides.add_slide(sl)
        slide.shapes.background.fill.solid()
        slide.shapes.background.fill.fore_color.rgb=bg
        return slide

    def slide_title(title):
        slide=blank()
        add_rect(slide,0.2,0.15,0.07,0.5,ORANGE)
        add_text(slide,title,0.35,0.1,12.5,0.6,size=22,bold=True,color=DARK)
        return slide

    dir_keys = [("Премиксы",nd["premixes"],pd.get("premixes",{})),
                ("Комбикорм",nd["compound"],pd.get("compound",{})),
                ("Кормовые добавки",nd["additives"],pd.get("additives",{}))]
    if nd["concentrates"].get("vol"):
        dir_keys.append(("Корм. концентраты",nd["concentrates"],pd.get("concentrates",{})))

    tVol=sum(float(d.get("vol") or 0) for _,d,_ in dir_keys)
    tVolP=sum(float(p.get("vol") or 0) for _,_,p in dir_keys)
    tRev=sum(float(d.get("revenue") or 0) for _,d,_ in dir_keys)
    tRevP=sum(float(p.get("revenue") or 0) for _,_,p in dir_keys)
    tGP=sum(float(d.get("gp") or 0) for _,d,_ in dir_keys)
    tGPP=sum(float(p.get("gp") or 0) for _,_,p in dir_keys)
    avgM=round(tGP/tRev*100,2) if tRev else 0
    avgMP=round(tGPP/tRevP*100,2) if tRevP else 0

    # Титул
    slide=blank(RGBColor(0xFF,0xFF,0xFF))
    add_rect(slide,W*0.56,0,W*0.44,H,ORANGE)
    add_text(slide,"Отчёт отдела кормления",0.5,2.0,7,0.8,size=32,bold=True,color=DARK)
    add_text(slide,f"за {month} 2026 г.",0.5,2.9,7,0.6,size=26,bold=True,color=ORANGE)
    add_text(slide,'ООО "КЛМ" | г. Минск',0.5,3.7,7,0.4,size=14,color=GRAY)
    add_text(slide,"ВМЕСТЕ СМОТРИМ В БУДУЩЕЕ!",W*0.58,H-0.5,W*0.4,0.4,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

    # Аналитика
    slide=slide_title(f"Аналитика по направлениям — {month}")
    cw=3.8; gp2=0.3; sx=0.4; cy=1.2; ch=2.2
    for i,(lbl,d,p) in enumerate(dir_keys[:3]):
        x=sx+i*(cw+gp2)
        add_rect(slide,x,cy,cw,ch,WHITE)
        add_text(slide,lbl,x+0.15,cy+0.1,cw-0.3,0.3,size=12,bold=True,color=ORANGE)
        add_text(slide,f"{fmt(d.get('vol',''),2)} т",x+0.15,cy+0.45,cw-0.3,0.5,size=22,bold=True,color=ORANGE)
        p2=pct(d.get("vol"),p.get("vol"))
        if p2 is not None:
            add_text(slide,dstr(p2),x+0.15,cy+1.0,cw-0.3,0.3,size=14,bold=True,color=dcol_pptx(p2))
        add_text(slide,f"vs {prev_month} ({fmt(p.get('vol',''),2)} т)",x+0.15,cy+1.4,cw-0.3,0.3,size=10,color=GRAY)

    # ИТОГО
    slide=slide_title(f"ИТОГО за {month} 2026")
    cards=[("Объём",f"{fmt(tVol,2)} т",pct(tVol,tVolP),f"{fmt(tVolP,2)} т"),
           ("Выручка",f"{fmt(tRev)} BYN",pct(tRev,tRevP),f"{fmt(tRevP)} BYN"),
           ("Валовая прибыль",f"{fmt(tGP)} BYN",pct(tGP,tGPP),f"{fmt(tGPP)} BYN"),
           ("Рентабельность",f"{avgM}%",round(avgM-avgMP,2),f"{avgMP}%")]
    for i,(lbl,val,dp,pv) in enumerate(cards):
        x=0.3+i*3.1; sy2=1.5; cw2=2.9; ch2=2.2
        add_rect(slide,x,sy2,cw2,ch2,WHITE)
        add_text(slide,lbl,x+0.15,sy2+0.1,cw2-0.3,0.3,size=11,color=GRAY)
        add_text(slide,val,x+0.15,sy2+0.45,cw2-0.3,0.55,size=18,bold=True,color=ORANGE)
        if dp is not None:
            add_text(slide,dstr(dp),x+0.15,sy2+1.05,cw2-0.3,0.3,size=12,bold=True,color=dcol_pptx(dp))
        add_text(slide,f"{prev_month}: {pv}",x+0.15,sy2+1.45,cw2-0.3,0.3,size=9,color=GRAY)

    # Таблицы сравнения
    def tbl_slide(title, labels, prevs, curs):
        slide=slide_title(title)
        rows=[["Направление",prev_month,month,"Динамика"]]
        for i,lbl in enumerate(labels):
            p2=pct(curs[i],prevs[i])
            rows.append([lbl,fmt(prevs[i]),fmt(curs[i]),dstr(p2)])
        tbl=slide.shapes.add_table(len(rows),4,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.55*len(rows))).table
        for r,row in enumerate(rows):
            for c,val in enumerate(row):
                cell=tbl.cell(r,c); cell.text=val
                p3=cell.text_frame.paragraphs[0]
                run=p3.runs[0] if p3.runs else p3.add_run()
                run.font.size=Pt(12 if r>0 else 11)
                run.font.bold=(r==0)
                run.font.color.rgb=DARK if r==0 else (GREEN if (c==3 and val.startswith("+")) else RED if (c==3 and val.startswith("-")) else DARK)
                cell.fill.solid()
                cell.fill.fore_color.rgb=RGBColor(0xFF,0xF8,0xF0) if r%2==0 else WHITE

    lbls=[k[0] for k in dir_keys]
    tbl_slide(f"Объёмы по направлениям (тонны)",      lbls,[float(p.get("vol") or 0) for _,_,p in dir_keys],[float(d.get("vol") or 0) for _,d,_ in dir_keys])
    tbl_slide(f"Валовая прибыль по направлениям (BYN)",lbls,[float(p.get("gp") or 0) for _,_,p in dir_keys],[float(d.get("gp") or 0) for _,d,_ in dir_keys])
    tbl_slide(f"Рентабельность по направлениям (%)",   lbls,[float(p.get("margin") or 0) for _,_,p in dir_keys],[float(d.get("margin") or 0) for _,d,_ in dir_keys])

    # Детали направлений
    for lbl,d,p in dir_keys:
        slide=slide_title(f"{lbl} — {month}")
        lines=[("Объём",fmt(d.get("vol",""),2)+" т",pct(d.get("vol"),p.get("vol")),fmt(p.get("vol",""),2)+" т"),
               ("Выручка",fmt(d.get("revenue",""))+" BYN (без НДС)",pct(d.get("revenue"),p.get("revenue")),fmt(p.get("revenue",""))+" BYN"),
               ("Валовая прибыль",fmt(d.get("gp",""))+" BYN",pct(d.get("gp"),p.get("gp")),fmt(p.get("gp",""))+" BYN"),
               ("Рентабельность",str(d.get("margin",""))+"%",round(float(d.get("margin") or 0)-float(p.get("margin") or 0),2),str(p.get("margin",""))+"%")]
        for i,(name,val,dp,pv) in enumerate(lines):
            y=1.1+i*0.78
            add_text(slide,f"• {name}:",0.4,y,3.5,0.55,size=15,color=DARK)
            add_text(slide,val,3.9,y,4.0,0.55,size=15,bold=True,color=DARK)
            if dp is not None:
                add_text(slide,dstr(dp),7.9,y,2.0,0.55,size=14,bold=True,color=dcol_pptx(dp))
            add_text(slide,f"было: {pv}",9.9,y,3.0,0.55,size=11,color=GRAY)

    # Сотрудники
    def person_slide(title, people, show_share=True):
        if not people: return
        slide=slide_title(title)
        tot=sum(float(m["gp"]) for m in people)
        mpr=4; cw3=min(3.0,12.5/min(len(people),mpr))
        for i,m in enumerate(people):
            row=i//mpr; col=i%mpr
            x=0.3+col*(cw3+0.15); y=1.1+row*2.3
            add_rect(slide,x,y,cw3,2.0,WHITE)
            add_text(slide,m["name"],x+0.1,y+0.1,cw3-0.2,0.3,size=11,bold=True,color=ORANGE)
            add_text(slide,fmt(m["gp"])+" BYN",x+0.1,y+0.45,cw3-0.2,0.45,size=16,bold=True,color=DARK)
            p2=pct(m["gp"],m["gp_prev"]) if m["gp_prev"]>0 else None
            if p2 is not None:
                add_text(slide,f"{dstr(p2)} к {prev_month}",x+0.1,y+0.95,cw3-0.2,0.28,size=11,color=dcol_pptx(p2))
            else:
                add_text(slide,"новый" if m["gp_prev"]==0 else "нет данных",x+0.1,y+0.95,cw3-0.2,0.28,size=10,color=GRAY)
            if show_share and tot>0:
                add_text(slide,f"{round(float(m['gp'])/tot*100,1)}% от ВП",x+0.1,y+1.28,cw3-0.2,0.25,size=9,color=GRAY)
        add_text(slide,f"ИТОГО: {fmt(tot)} BYN",0.3,H-0.6,8,0.4,size=13,bold=True,color=DARK)

    person_slide("Менеджеры отдела кормления",nd["dept_managers"],True)
    person_slide("Консультанты",nd["consultants"],True)
    person_slide("Менеджеры других отделов",nd["other_managers"],False)

    # ВП отдела
    slide=slide_title(f"ВП отдела кормления: {month} vs {prev_month}")
    dept_n=nd["dept"]; dept_p=pd.get("dept",{})
    for i,(lbl,cv,pv) in enumerate([("ВП по приходу",dept_n["gp_arrival"],dept_p.get("gp_arrival",0)),("Транспорт",dept_n["transport"],dept_p.get("transport",0)),("ВП к расчёту ЗП",dept_n["gp_salary"],dept_p.get("gp_salary",0))]):
        x=0.4+i*4.3; y=1.5; cw4=4.0; ch4=2.4
        add_rect(slide,x,y,cw4,ch4,WHITE)
        add_text(slide,lbl,x+0.15,y+0.1,cw4-0.3,0.3,size=11,color=GRAY)
        add_text(slide,fmt(cv)+" BYN",x+0.15,y+0.5,cw4-0.3,0.55,size=20,bold=True,color=ORANGE)
        p2=pct(cv,pv)
        if p2: add_text(slide,dstr(p2),x+0.15,y+1.1,cw4-0.3,0.3,size=13,bold=True,color=dcol_pptx(p2))
        add_text(slide,f"{prev_month}: {fmt(pv)} BYN",x+0.15,y+1.55,cw4-0.3,0.3,size=10,color=GRAY)

    # ВСЕГО
    slide=slide_title("ВСЕГО по направлению кормления")
    dt_n=nd["dir_total"]; dt_p=pd.get("dir_total",{})
    for i,(lbl,cv,pv) in enumerate([("Валовая прибыль",dt_n["gp"],dt_p.get("gp",0)),("Транспортные затраты",dt_n["transport"],dt_p.get("transport",0))]):
        x=0.4+i*6.5; y=1.8; cw5=6.0; ch5=2.4
        add_rect(slide,x,y,cw5,ch5,WHITE)
        add_text(slide,lbl,x+0.15,y+0.1,cw5-0.3,0.3,size=12,color=GRAY)
        add_text(slide,fmt(cv)+" BYN",x+0.15,y+0.5,cw5-0.3,0.6,size=26,bold=True,color=ORANGE)
        p2=pct(cv,pv)
        if p2: add_text(slide,dstr(p2),x+0.15,y+1.2,cw5-0.3,0.35,size=14,bold=True,color=dcol_pptx(p2))
        add_text(slide,f"{prev_month}: {fmt(pv)} BYN",x+0.15,y+1.65,cw5-0.3,0.3,size=10,color=GRAY)

    # Финал
    slide=blank(RGBColor(0xFF,0xFF,0xFF))
    add_rect(slide,W*0.55,0,W*0.45,H,ORANGE)
    add_text(slide,"Спасибо за внимание!",0.5,2.8,7,1.0,size=36,bold=True,color=DARK)
    add_text(slide,"ВМЕСТЕ СМОТРИМ В БУДУЩЕЕ!",W*0.57,H-0.55,W*0.42,0.4,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(slide,'ООО "КЛМ" | г. Минск',W*0.57,H-0.3,W*0.42,0.28,size=11,color=WHITE,align=PP_ALIGN.CENTER)

    buf=io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf.getvalue()

# ============================================================
# ГЕНЕРАЦИЯ PDF
# ============================================================
def build_pdf(nd, pd, month, prev_month):
    buf=io.BytesIO()
    doc=SimpleDocTemplate(buf,pagesize=A4,leftMargin=1.5*cm,rightMargin=1.5*cm,topMargin=1.5*cm,bottomMargin=1.5*cm)
    styles=getSampleStyleSheet()

    def S(name,**kw):
        return ParagraphStyle(name,parent=styles["Normal"],
            fontName=kw.pop("font","DejaVu"),fontSize=kw.pop("size",11),
            textColor=kw.pop("color",DARK_PDF),leading=kw.pop("leading",16),**kw)

    sH1=S("h1",font="DejaVuBold",size=22,color=ORANGE_PDF,leading=28)
    sH2=S("h2",font="DejaVuBold",size=14,color=DARK_PDF,leading=20,spaceBefore=10)
    sSmall=S("sm",size=9,color=GRAY_PDF,leading=13)
    sCenter=S("sc",size=14,font="DejaVuBold",alignment=TA_CENTER,color=DARK_PDF,leading=20)
    sOrgCenter=S("oc",size=14,font="DejaVuBold",alignment=TA_CENTER,color=ORANGE_PDF,leading=20)
    sGrayCenter=S("gc",size=11,alignment=TA_CENTER,color=GRAY_PDF,leading=16)

    TS=lambda: TableStyle([
        ("BACKGROUND",(0,0),(-1,0),ORANGE_PDF),
        ("TEXTCOLOR",(0,0),(-1,0),colors.white),
        ("FONTNAME",(0,0),(-1,0),"DejaVuBold"),
        ("FONTNAME",(0,1),(-1,-1),"DejaVu"),
        ("FONTSIZE",(0,0),(-1,-1),9),
        ("ROWBACKGROUNDS",(0,1),(-1,-1),[CREAM_PDF,colors.white]),
        ("GRID",(0,0),(-1,-1),0.5,colors.lightgrey),
        ("ALIGN",(1,0),(-1,-1),"CENTER"),
        ("VALIGN",(0,0),(-1,-1),"MIDDLE"),
        ("TOPPADDING",(0,0),(-1,-1),4),
        ("BOTTOMPADDING",(0,0),(-1,-1),4),
    ])

    def section(title):
        story.append(Spacer(1,0.3*cm))
        story.append(HRFlowable(width="100%",thickness=2,color=ORANGE_PDF))
        story.append(Paragraph(title,sH2))
        story.append(Spacer(1,0.2*cm))

    def make_table(rows,cws):
        t=Table(rows,colWidths=cws); t.setStyle(TS()); return t

    story=[]
    dir_keys=[("Премиксы",nd["premixes"],pd.get("premixes",{})),
              ("Комбикорм",nd["compound"],pd.get("compound",{})),
              ("Кормовые добавки",nd["additives"],pd.get("additives",{}))]
    if nd["concentrates"].get("vol"):
        dir_keys.append(("Корм. концентраты",nd["concentrates"],pd.get("concentrates",{})))

    tVol=sum(float(d.get("vol") or 0) for _,d,_ in dir_keys)
    tVolP=sum(float(p.get("vol") or 0) for _,_,p in dir_keys)
    tRev=sum(float(d.get("revenue") or 0) for _,d,_ in dir_keys)
    tRevP=sum(float(p.get("revenue") or 0) for _,_,p in dir_keys)
    tGP=sum(float(d.get("gp") or 0) for _,d,_ in dir_keys)
    tGPP=sum(float(p.get("gp") or 0) for _,_,p in dir_keys)
    avgM=round(tGP/tRev*100,2) if tRev else 0
    avgMP=round(tGPP/tRevP*100,2) if tRevP else 0

    # Титул
    story.append(Spacer(1,0.5*cm))
    story.append(Paragraph("Отчёт отдела кормления",sH1))
    story.append(Paragraph(f"за {month} 2026 г.",S("m2",font="DejaVuBold",size=16,color=DARK_PDF,leading=22)))
    story.append(Paragraph('ООО "КЛМ" | г. Минск',sSmall))
    story.append(Spacer(1,0.4*cm))

    # Аналитика
    section(f"Аналитика по направлениям — {month}")
    rows=[["Направление",f"Объём {month}",f"Объём {prev_month}","Динамика"]]
    for lbl,d,p in dir_keys:
        rows.append([lbl,fmt(d.get("vol",""),2)+" т",fmt(p.get("vol",""),2)+" т",dstr(pct(d.get("vol"),p.get("vol")))])
    story.append(make_table(rows,[5*cm,4*cm,4*cm,3*cm]))

    # ИТОГО
    section(f"ИТОГО за {month} 2026")
    rows=[["Показатель",month,prev_month,"Динамика"],
          ["Объём (т)",fmt(tVol,2),fmt(tVolP,2),dstr(pct(tVol,tVolP))],
          ["Выручка (BYN)",fmt(tRev),fmt(tRevP),dstr(pct(tRev,tRevP))],
          ["Валовая прибыль (BYN)",fmt(tGP),fmt(tGPP),dstr(pct(tGP,tGPP))],
          ["Рентабельность (%)",str(avgM),str(avgMP),dstr(round(avgM-avgMP,2))]]
    story.append(make_table(rows,[6*cm,3.5*cm,3.5*cm,3*cm]))

    # Детали
    for lbl,d,p in dir_keys:
        section(f"{lbl} — {month}")
        rows=[["Показатель",month,prev_month,"Динамика"],
              ["Объём (т)",fmt(d.get("vol",""),2),fmt(p.get("vol",""),2),dstr(pct(d.get("vol"),p.get("vol")))],
              ["Выручка (BYN)",fmt(d.get("revenue","")),fmt(p.get("revenue","")),dstr(pct(d.get("revenue"),p.get("revenue")))],
              ["Валовая прибыль",fmt(d.get("gp","")),fmt(p.get("gp","")),dstr(pct(d.get("gp"),p.get("gp")))],
              ["Рентабельность (%)",str(d.get("margin","")),str(p.get("margin","")),dstr(round(float(d.get("margin") or 0)-float(p.get("margin") or 0),2))]]
        story.append(make_table(rows,[6*cm,3.5*cm,3.5*cm,3*cm]))

    # Сотрудники
    def people_tbl(title,people,show_share=True):
        if not people: return
        section(title)
        tot=sum(float(m["gp"]) for m in people)
        hdr=["Имя",f"ВП {month}",f"ВП {prev_month}","Динамика"]
        if show_share: hdr.append("Доля")
        rows=[hdr]
        for m in people:
            p2=pct(m["gp"],m["gp_prev"]) if m["gp_prev"]>0 else None
            row=[m["name"],fmt(m["gp"])+" BYN",fmt(m["gp_prev"])+" BYN" if m["gp_prev"] else "—",dstr(p2)]
            if show_share: row.append(f"{round(float(m['gp'])/tot*100,1)}%" if tot>0 else "")
            rows.append(row)
        tot_row=["ИТОГО",fmt(tot)+" BYN","",""]
        if show_share: tot_row.append("")
        rows.append(tot_row)
        cws=[5*cm,3*cm,3*cm,2.5*cm]
        if show_share: cws.append(2*cm)
        ts=TS()
        ts.add("FONTNAME",(0,-1),(-1,-1),"DejaVuBold")
        ts.add("BACKGROUND",(0,-1),(-1,-1),CREAM_PDF)
        t=Table(rows,colWidths=cws); t.setStyle(ts)
        story.append(t)

    people_tbl("Менеджеры отдела кормления",nd["dept_managers"],True)
    people_tbl("Консультанты",nd["consultants"],True)
    people_tbl("Менеджеры других отделов",nd["other_managers"],False)

    # ВП отдела
    section(f"ВП отдела кормления: {month} vs {prev_month}")
    dept_n=nd["dept"]; dept_p=pd.get("dept",{})
    rows=[["Показатель",month,prev_month,"Динамика"],
          ["ВП по приходу",fmt(dept_n["gp_arrival"])+" BYN",fmt(dept_p.get("gp_arrival",0))+" BYN",dstr(pct(dept_n["gp_arrival"],dept_p.get("gp_arrival",0)))],
          ["Транспорт",fmt(dept_n["transport"])+" BYN",fmt(dept_p.get("transport",0))+" BYN",dstr(pct(dept_n["transport"],dept_p.get("transport",0)))],
          ["ВП к расчёту ЗП",fmt(dept_n["gp_salary"])+" BYN",fmt(dept_p.get("gp_salary",0))+" BYN",dstr(pct(dept_n["gp_salary"],dept_p.get("gp_salary",0)))]]
    story.append(make_table(rows,[5*cm,4*cm,4*cm,3*cm]))

    # ВСЕГО
    section("ВСЕГО по направлению кормления")
    dt_n=nd["dir_total"]; dt_p=pd.get("dir_total",{})
    rows=[["Показатель",month,prev_month,"Динамика"],
          ["Валовая прибыль",fmt(dt_n["gp"])+" BYN",fmt(dt_p.get("gp",0))+" BYN",dstr(pct(dt_n["gp"],dt_p.get("gp",0)))],
          ["Транспортные затраты",fmt(dt_n["transport"])+" BYN",fmt(dt_p.get("transport",0))+" BYN",dstr(pct(dt_n["transport"],dt_p.get("transport",0)))]]
    story.append(make_table(rows,[5*cm,4*cm,4*cm,3*cm]))

    # Финал
    story.append(Spacer(1,1*cm))
    story.append(HRFlowable(width="100%",thickness=2,color=ORANGE_PDF))
    story.append(Spacer(1,0.5*cm))
    story.append(Paragraph("Спасибо за внимание!",sCenter))
    story.append(Paragraph("ВМЕСТЕ СМОТРИМ В БУДУЩЕЕ!",sOrgCenter))
    story.append(Paragraph('ООО "КЛМ" | г. Минск',sGrayCenter))

    doc.build(story)
    buf.seek(0)
    return buf.getvalue()

# ============================================================
# STREAMLIT UI
# ============================================================
st.set_page_config(page_title='Отчёт КЛМ', page_icon='🐄', layout='centered')
st.markdown(f"""
<div style='background:#E87722;padding:20px 28px;border-radius:12px;margin-bottom:20px'>
  <h1 style='color:white;margin:0;font-size:26px'>🐄 Генератор отчётов</h1>
  <p style='color:rgba(255,255,255,0.85);margin:6px 0 0'>ООО "КЛМ" | Отдел кормления</p>
</div>
""", unsafe_allow_html=True)

col1, col2 = st.columns(2)
with col1:
    st.subheader("📊 Новый месяц")
    new_file = st.file_uploader("Excel файл нового месяца", type=["xlsx","xls"], key="new")
with col2:
    st.subheader("📋 Прошлый месяц")
    prev_file = st.file_uploader("Excel файл прошлого месяца", type=["xlsx","xls"], key="prev")

col3, col4 = st.columns(2)
with col3:
    month = st.text_input("Отчётный месяц", value="апрель", placeholder="напр.: апрель")
with col4:
    prev_month = st.text_input("Прошлый месяц", value="март", placeholder="напр.: март")

if st.button("⚙️ Сгенерировать отчёт", type="primary", use_container_width=True):
    if not new_file or not prev_file:
        st.error("Загрузите оба Excel файла!")
    elif not month or not prev_month:
        st.error("Введите названия месяцев!")
    else:
        with st.spinner("⏳ Генерация PPTX и PDF..."):
            try:
                nd = parse_excel(new_file.read())
                pd_data = parse_excel(prev_file.read())

                # Прошлые данные сотрудников
                for person in nd["dept_managers"]:
                    found=next((p for p in pd_data["dept_managers"] if p["name"].split()[0]==person["name"].split()[0]),None)
                    if found: person["gp_prev"]=found["gp"]
                for person in nd["consultants"]:
                    found=next((p for p in pd_data["consultants"] if p["name"].split()[0]==person["name"].split()[0]),None)
                    if found: person["gp_prev"]=found["gp"]
                for person in nd["other_managers"]:
                    found=next((p for p in pd_data["other_managers"] if p["name"].split()[0]==person["name"].split()[0]),None)
                    if found: person["gp_prev"]=found["gp"]

                pd_merged = {
                    "premixes":pd_data["premixes"],"compound":pd_data["compound"],
                    "additives":pd_data["additives"],"concentrates":pd_data["concentrates"],
                    "dept":pd_data["dept"],"dir_total":pd_data["dir_total"],
                }

                pptx_bytes = build_pptx(nd, pd_merged, month, prev_month)
                pdf_bytes  = build_pdf(nd, pd_merged, month, prev_month)

                st.success("✅ Готово! Скачайте файлы:")
                c1,c2=st.columns(2)
                with c1:
                    st.download_button(
                        label="📥 Скачать PPTX",
                        data=pptx_bytes,
                        file_name=f"Отчёт_кормление_{month}_2026.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                with c2:
                    st.download_button(
                        label="📄 Скачать PDF",
                        data=pdf_bytes,
                        file_name=f"Отчёт_кормление_{month}_2026.pdf",
                        mime="application/pdf",
                        use_container_width=True
                    )
            except Exception as e:
                st.error(f"Ошибка: {e}")
                st.exception(e)
